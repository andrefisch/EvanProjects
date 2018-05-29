from PIL import Image
from pptx import Presentation
import io
import openpyxl
import os
import pygame
import time


'''
- loop through slides
- store information in an array
- for each shape in the slide record [row, text, source]
'''
def extract_info_from_slides():
    pptxs = []
    files = [x for x in os.listdir() if x.endswith(".pptx")]
    count = 1
    for eachfile in files:
        prs = Presentation(eachfile)
        print("----------------------")
        print(eachfile)
        print("----------------------")
        for slide in prs.slides:
            for shape in slide.shapes:
                info = []
                if hasattr(shape, "image"):
                    count = count + 1
                    image = Image.open(io.BytesIO(shape.image.blob)).convert('RGB')
                    name = str(count) + "image.jpg"
                    print(name)
                    image.save(name)
                    info = [count, name]
                    pptxs.append(info)
                    print(str(count) + ": " + "FOUND AN IMAGE! " + str(shape.image.size))
                if hasattr(shape, "text"):
                    count = count + 1
                    info = [count, shape.text, eachfile]
                    pptxs.append(info)
                    print(str(count) + ": " + shape.text)
    return pptxs

def slides_to_excel():
    # Create a new excel file
    out = openpyxl.Workbook()
    # Open the worksheet we want to edit
    outsheet = out.create_sheet("slides")

    # if 'sheet' appears randomly we can delete it
    rm = out.get_sheet_by_name('Sheet')
    out.remove_sheet(rm)

    #################
    # DO STUFF HERE #
    #################
    TEXT   = "A"
    SOURCE = "B"

    outsheet[TEXT   + '1'].value = "Text"
    outsheet[SOURCE + '1'].value = "Source"
    
    pptxs = extract_info_from_slides()
    for i in range(0, len(pptxs)):
        slide = pptxs[i]
        if len(slide) == 2:
            # img = openpyxl.drawing.image.Image(slide[1])
            # THIS LINE WORKS, PROBLEM IS ABOVE
            print(slide[1])
            # img = Image.open('3image.jpg')
            # img = Image(slide[1])
            img = openpyxl.drawing.image.Image(slide[1])
            img.anchor(outsheet.cell(TEXT + str(slide[0])))
            outsheet.add_image(img)
        elif len(slide) == 3:
            outsheet[TEXT   + str(slide[0])].value = slide[1]
            outsheet[SOURCE + str(slide[0])].value = slide[2]

    os.remove("*.jpg")

    # Save the file
    out.save("newFile.xlsx")

    # LMK when the script is done
    pygame.init()
    pygame.mixer.music.load('/home/andrefisch/python/evan/note.mp3')
    pygame.mixer.music.play()
    time.sleep(3)
    pygame.mixer.music.stop()

# extract_info_from_slides()
slides_to_excel()
